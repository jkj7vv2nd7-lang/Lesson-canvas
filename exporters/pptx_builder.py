"""
「板書用スライド構成」のMarkdown表を、実際のpptxファイルに変換する。

想定する表構造（core/prompts.py の ARTIFACT_FORMATS["slide_outline"]）:
  スライド番号 | 見出し | 本文要点 | 話者ノート

AI生成テキストの表記ゆれ（列の過不足、区切り文字の違いなど）に対して
できるだけ寛容に解析する。
"""

from __future__ import annotations

import io
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from core.md_parse import extract_tables, split_into_sections

TITLE_COLOR = RGBColor(0x2B, 0x58, 0x76)
TEXT_COLOR = RGBColor(0x22, 0x22, 0x22)
BG_COLOR = RGBColor(0xFF, 0xFF, 0xFF)


def _split_bullets(text: str) -> list[str]:
    """本文要点セルを箇条書き単位に分割する（・や句点区切りに対応）。"""
    if not text:
        return []
    parts = re.split(r"[・]|(?<=[。])\s*", text)
    return [p.strip() for p in parts if p.strip()]


def _set_background_white(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def build_slide_outline_pptx(
    *,
    deck_title: str,
    content_text: str,
    footer_note: str = "",
) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]  # 白紙レイアウト

    # ---- 表紙 ----
    cover = prs.slides.add_slide(blank_layout)
    _set_background_white(cover)
    title_box = cover.shapes.add_textbox(Inches(1), Inches(2.8), Inches(11.3), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = deck_title
    run.font.size = Pt(40)
    run.font.bold = True
    run.font.color.rgb = TITLE_COLOR

    # ---- 表データの抽出 ----
    sections = split_into_sections(content_text)
    rows: list[list[str]] = []
    for sec in sections:
        tables = extract_tables(sec.body_lines)
        if tables:
            rows = tables[0]
            break
    if not rows:
        # セクション見出しがない場合、全文から表を探す
        tables = extract_tables(content_text.split("\n"))
        if tables:
            rows = tables[0]

    if len(rows) >= 2:
        header = [h.strip() for h in rows[0]]

        def col_index(*keywords: str) -> int | None:
            for i, h in enumerate(header):
                if any(k in h for k in keywords):
                    return i
            return None

        idx_heading = col_index("見出し", "タイトル")
        idx_body = col_index("本文", "要点")
        idx_notes = col_index("話者", "ノート", "発問")

        for row in rows[1:]:
            slide = prs.slides.add_slide(blank_layout)
            _set_background_white(slide)

            heading_text = row[idx_heading] if idx_heading is not None and idx_heading < len(row) else ""
            body_text = row[idx_body] if idx_body is not None and idx_body < len(row) else ""
            notes_text = row[idx_notes] if idx_notes is not None and idx_notes < len(row) else ""

            # 見出し
            head_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.9), Inches(1.1))
            htf = head_box.text_frame
            htf.word_wrap = True
            hp = htf.paragraphs[0]
            hrun = hp.add_run()
            hrun.text = heading_text or f"（見出し未設定）"
            hrun.font.size = Pt(32)
            hrun.font.bold = True
            hrun.font.color.rgb = TITLE_COLOR

            # 本文（箇条書き、最大3点まで。多すぎる場合は分割の目安として警告的に全件出す）
            body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.9), Inches(11.5), Inches(4.8))
            btf = body_box.text_frame
            btf.word_wrap = True
            bullets = _split_bullets(body_text)
            if not bullets:
                bullets = [body_text] if body_text else []
            for i, bullet in enumerate(bullets):
                bp = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
                bp.level = 0
                brun = bp.add_run()
                brun.text = f"・{bullet}"
                brun.font.size = Pt(22)
                brun.font.color.rgb = TEXT_COLOR
                bp.space_after = Pt(14)

            # 話者ノート
            if notes_text:
                slide.notes_slide.notes_text_frame.text = notes_text

    # ---- 注記スライド ----
    if footer_note:
        note_slide = prs.slides.add_slide(blank_layout)
        _set_background_white(note_slide)
        note_box = note_slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.3), Inches(1.5))
        ntf = note_box.text_frame
        ntf.word_wrap = True
        np_ = ntf.paragraphs[0]
        nrun = np_.add_run()
        nrun.text = footer_note
        nrun.font.size = Pt(14)
        nrun.font.color.rgb = RGBColor(0x70, 0x70, 0x70)
        nrun.font.italic = True

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
